from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "The AI Engineering Journey"
subtitle.text = "Insights on Antigravity, Career Roadmaps, and the Confide Project\nPrepared by Antigravity"

# Slide 2: How is Antigravity Built?
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "How is an AI Agent (Antigravity) Built?"
tf = body_shape.text_frame
tf.text = "Built in four massive layers:"

p = tf.add_paragraph()
p.text = "1. The Brain: A core LLM (like Gemini) trained on millions of repositories (Python/C++)."
p.level = 1

p = tf.add_paragraph()
p.text = "2. The Hands: Tool-calling capabilities that let the AI interact with files and terminals."
p.level = 1

p = tf.add_paragraph()
p.text = "3. The Loop: An 'Agentic ReAct Loop' (Reason + Act) allowing autonomous execution."
p.level = 1

p = tf.add_paragraph()
p.text = "4. The Sandbox: A secure environment ensuring safe execution of commands."
p.level = 1

# Slide 3: AI Building AI
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Is AI Completely Hand-Coded?"
tf = body_shape.text_frame
tf.text = "The Short Answer: No!"

p = tf.add_paragraph()
p.text = "The Early Days (2017): 100% human-coded (grueling 15-hour days by researchers)."
p.level = 1

p = tf.add_paragraph()
p.text = "Today: AI Engineers use AI to build better AI!"
p.level = 1

p = tf.add_paragraph()
p.text = "Humans are now 'Architects', while AI handles the repetitive 'typing' and boilerplate."
p.level = 1

# Slide 4: Confide Project Review
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Review: Your 'Confide AI' Project"
tf = body_shape.text_frame
tf.text = "An exceptionally mature project for a 2nd-year B.Tech student."

p = tf.add_paragraph()
p.text = "Architecture: Modern Full-Stack (React/Vite + Node.js/Express)."
p.level = 1

p = tf.add_paragraph()
p.text = "Maturity: Features like 'Crisis Protocols' show real Product Management thinking."
p.level = 1

p = tf.add_paragraph()
p.text = "Flexibility: Integrates multiple AI providers (Groq, Gemini, OpenAI) + Vision support."
p.level = 1

p = tf.add_paragraph()
p.text = "Next Steps: Add a database (MongoDB/Supabase) and User Authentication."
p.level = 1

# Slide 5: The Career Roadmap
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "The Golden Career Roadmap"
tf = body_shape.text_frame
tf.text = "How to maximize ROI in your CS degree:"

p = tf.add_paragraph()
p.text = "1. Be a Full-Stack AI Engineer: Combine Web Dev (React) with AI integrations."
p.level = 1

p = tf.add_paragraph()
p.text = "2. Master RAG & Agents: Companies pay top dollar for engineers who can make models useful."
p.level = 1

p = tf.add_paragraph()
p.text = "3. End-to-End Deployment: Learn Docker and Cloud (AWS/Vercel) to deploy your apps."
p.level = 1

p = tf.add_paragraph()
p.text = "Avoid 'Vibe Coding' trap: Use AI as a tutor, not a replacement for understanding."
p.level = 1

# Save the presentation
prs.save("Antigravity_Chat_Summary.pptx")
print("Presentation saved successfully as Antigravity_Chat_Summary.pptx")
